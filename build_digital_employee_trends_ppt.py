from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os
import zipfile


SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.8)
CONTENT_W = SW - Inches(1.6)
SOURCE_Y = Inches(7.03)

NAVY = RGBColor(0x05, 0x1C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
MED = RGBColor(0x66, 0x66, 0x66)
LINE = RGBColor(0xCC, 0xCC, 0xCC)
BG = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_BLUE = RGBColor(0xE9, 0xF6, 0xFD)
LIGHT_GREEN = RGBColor(0xE9, 0xF5, 0xEE)
LIGHT_ORANGE = RGBColor(0xFD, 0xF2, 0xE9)
LIGHT_RED = RGBColor(0xFB, 0xEB, 0xEB)
BLUE = RGBColor(0x00, 0xA9, 0xF4)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)

TITLE_SIZE = Pt(22)
SUB_SIZE = Pt(18)
BODY_SIZE = Pt(14)
SMALL_SIZE = Pt(9)
BIG_SIZE = Pt(30)

COMMON_SOURCE = 'Source: Synthesis based on McKinsey State of AI (2025), McKinsey Tech Trends Outlook (2025), and Microsoft Work Trend Index (2025)'


def clean_shape(shape):
    style = shape._element.find(qn('p:style'))
    if style is not None:
        shape._element.remove(style)


def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def add_text(slide, left, top, width, height, text, font_size=BODY_SIZE,
             font_name='Arial', font_color=DARK, bold=False,
             alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             paragraph_space=Pt(6), line_multiple=1.28):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    body_pr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    body_pr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        body_pr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = paragraph_space if idx > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * line_multiple)
        for run in p.runs:
            set_ea_font(run)
    return box


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    clean_shape(shape)
    return shape


def add_hline(slide, left, top, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, left, top, length, h, color)


def add_oval(slide, left, top, label, size=Inches(0.46), bg=NAVY, fg=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    clean_shape(shape)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = BODY_SIZE
    p.font.name = 'Arial'
    p.font.color.rgb = fg
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    body_pr = tf._txBody.find(qn('a:bodyPr'))
    body_pr.set('anchor', 'ctr')
    for run in p.runs:
        set_ea_font(run)
    return shape


def add_action_title(slide, text):
    add_text(slide, LM, Inches(0.15), Inches(11.7), Inches(0.9), text,
             font_size=TITLE_SIZE, font_name='Georgia', font_color=BLACK,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, LM, Inches(1.03), Inches(11.7), BLACK, Pt(0.5))


def add_source(slide, text, page_num=None):
    add_text(slide, Inches(0.8), SOURCE_Y, Inches(10.7), Inches(0.22), text,
             font_size=SMALL_SIZE, font_color=MED)
    if page_num is not None:
        add_text(slide, Inches(12.0), SOURCE_Y, Inches(0.55), Inches(0.22), str(page_num),
                 font_size=SMALL_SIZE, font_color=MED, alignment=PP_ALIGN.RIGHT)


def full_cleanup(path):
    tmp = path + '.tmp'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith('.xml'):
                root = etree.fromstring(data)
                for style in root.findall(f'.//{{{ns_p}}}style'):
                    style.getparent().remove(style)
                if 'theme' in item.filename.lower():
                    for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                        for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                            el.getparent().remove(el)
                data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes')
            zout.writestr(item, data)
    os.replace(tmp, path)


def bullet_block(slide, x, y, w, items, h, color=DARK, size=BODY_SIZE, spacing=Pt(10), align=PP_ALIGN.LEFT):
    add_text(slide, x, y, w, h, [f'• {item}' for item in items], font_size=size,
             font_color=color, paragraph_space=spacing, alignment=align)


def card(slide, x, y, w, h, title, body_lines, accent=NAVY, fill=BG, title_color=WHITE):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, w, Inches(0.55), accent)
    add_text(slide, x + Inches(0.12), y, w - Inches(0.24), Inches(0.55), title,
             font_size=BODY_SIZE, font_color=title_color, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.18), y + Inches(0.75), w - Inches(0.36), h - Inches(0.95),
             body_lines, font_size=Pt(12), font_color=DARK, paragraph_space=Pt(8))


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
    add_text(s, Inches(1.0), Inches(1.85), Inches(11.3), Inches(1.25),
             'Digital Employees', font_size=Pt(30), font_name='Georgia',
             font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.75),
             'Where the market is heading — and what enterprises should build next',
             font_size=Pt(20), font_color=DARK, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.35), Inches(4.05), Inches(10.7), Inches(1.2), [
             'A McKinsey-style perspective on how AI agents are evolving',
             'from copilots into governed digital labor'
             ], font_size=BODY_SIZE, font_color=MED, alignment=PP_ALIGN.CENTER, paragraph_space=Pt(8))
    add_text(s, Inches(1.0), Inches(5.05), Inches(11.3), Inches(0.4),
             'March 2026', font_size=BODY_SIZE, font_color=MED, alignment=PP_ALIGN.CENTER)
    add_hline(s, Inches(1.0), Inches(6.75), Inches(3.2), NAVY, Pt(2))
    add_text(s, Inches(1.0), Inches(6.92), Inches(11.0), Inches(0.22),
             'Generated with python-pptx in a McKinsey-style layout',
             font_size=SMALL_SIZE, font_color=MED)


def slide_exec_summary(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'Executive summary | Digital employees are moving from assistance to accountable execution')
    add_rect(s, LM, Inches(1.42), CONTENT_W, Inches(0.9), NAVY)
    add_text(s, Inches(1.02), Inches(1.52), Inches(10.9), Inches(0.65),
             'The winners will not be the firms that buy the most AI tools, but the ones that redesign workflows, control layers, and talent models around digital labor.',
             font_size=SUB_SIZE, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    cols = [
        ('1', 'Outcome ownership', 'The market is shifting from chat-based help toward agents that complete measurable pieces of work end to end.'),
        ('2', 'Role specialization', 'Value is moving from generic assistants to role-based digital employees in sales, service, operations, finance, and HR.'),
        ('3', 'System action', 'Competitive advantage comes from agents that can read context, use enterprise tools, and trigger actions inside workflows.'),
        ('4', 'Governed scale', 'As pilots scale, differentiation shifts from model novelty to integration, evaluation, security, and auditability.'),
    ]
    cw = Inches(2.72)
    gap = Inches(0.2)
    for i, (idx, title, desc) in enumerate(cols):
        x = LM + (cw + gap) * i
        add_rect(s, x, Inches(2.7), cw, Inches(3.15), BG)
        add_oval(s, x + Inches(0.12), Inches(2.82), idx)
        add_text(s, x + Inches(0.66), Inches(2.8), cw - Inches(0.82), Inches(0.42),
                 title, font_size=BODY_SIZE, font_color=NAVY, bold=True)
        add_hline(s, x + Inches(0.15), Inches(3.32), cw - Inches(0.3), LINE)
        add_text(s, x + Inches(0.15), Inches(3.5), cw - Inches(0.3), Inches(2.1),
                 desc, font_size=Pt(12), font_color=DARK)
    add_source(s, COMMON_SOURCE, num)


def slide_why_now(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'Why now | Four forces are pushing digital employees from concept to enterprise agenda')
    add_text(s, LM, Inches(1.48), Inches(11.6), Inches(0.4),
             'The acceleration is not driven by one model breakthrough alone; it is the result of simultaneous progress in capability, cost, access, and executive demand.',
             font_size=SUB_SIZE, font_color=NAVY, bold=True)
    reasons = [
        ('Better reasoning and tool use', ['Newer models plan, call tools, and manage multi-step tasks more reliably than simple chat assistants.'], BLUE, LIGHT_BLUE),
        ('More enterprise surface area', ['APIs, knowledge bases, SaaS connectors, and event streams make it easier for agents to act inside real workflows.'], GREEN, LIGHT_GREEN),
        ('Lower orchestration friction', ['LLM access, prompt tooling, evaluation stacks, and workflow builders have become easier to test and deploy.'], ORANGE, LIGHT_ORANGE),
        ('Capacity pressure at the top', ['Leaders want productivity gains, faster response times, and digital leverage without linear headcount growth.'], RED, LIGHT_RED),
    ]
    xs = [Inches(0.8), Inches(6.72)]
    ys = [Inches(2.0), Inches(4.2)]
    idx = 0
    for y in ys:
        for x in xs:
            title, body, accent, fill = reasons[idx]
            add_rect(s, x, y, Inches(5.8), Inches(1.75), fill)
            add_rect(s, x, y, Inches(0.14), Inches(1.75), accent)
            add_text(s, x + Inches(0.28), y + Inches(0.18), Inches(5.2), Inches(0.34),
                     title, font_size=SUB_SIZE, font_color=NAVY, bold=True)
            add_text(s, x + Inches(0.28), y + Inches(0.68), Inches(5.2), Inches(0.78),
                     body, font_size=Pt(12), font_color=DARK)
            idx += 1
    add_rect(s, Inches(0.8), Inches(6.33), Inches(11.7), Inches(0.38), BG)
    add_text(s, Inches(1.0), Inches(6.38), Inches(11.2), Inches(0.22),
             'Implication: most enterprises no longer need to ask whether digital employees are possible; they need to decide where they should be trusted first.',
             font_size=Pt(12), font_color=NAVY, bold=True)
    add_source(s, COMMON_SOURCE, num)


def slide_trends_a(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'Trend cluster A | The market is shifting toward more autonomous and more specialized digital workers')
    card(s, Inches(0.8), Inches(1.55), Inches(3.75), Inches(4.9), '1. From copilots to digital workers', [
        'Assistive chat remains useful, but spending is moving to systems that own a defined output.',
        'Examples include ticket resolution, proposal drafting, collections follow-up, and exception handling.',
        'The key buying question becomes: what work can the agent close, not just support?'
    ], accent=BLUE, fill=LIGHT_BLUE)
    card(s, Inches(4.79), Inches(1.55), Inches(3.75), Inches(4.9), '2. From general tools to role-specific agents', [
        'Boards care less about generic AI access and more about measurable agents for real business roles.',
        'Role context, policy rules, knowledge sources, and KPI alignment are becoming the new moat.',
        'Verticalization is where adoption quality and willingness to pay both improve.'
    ], accent=GREEN, fill=LIGHT_GREEN)
    card(s, Inches(8.78), Inches(1.55), Inches(3.75), Inches(4.9), '3. From one agent to agent teams', [
        'Complex work increasingly requires planners, researchers, reviewers, and executors to collaborate.',
        'This creates a digital team design problem, not just a prompt design problem.',
        'Multi-agent patterns will matter most where work has handoffs, controls, and iteration.'
    ], accent=ORANGE, fill=LIGHT_ORANGE)
    add_source(s, COMMON_SOURCE, num)


def slide_trends_b(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'Trend cluster B | Embedding, multimodality, and governance are defining the next wave of scale')
    card(s, Inches(0.8), Inches(1.55), Inches(3.75), Inches(4.9), '4. From text-only to multimodal work', [
        'Digital employees are increasingly expected to process documents, spreadsheets, screenshots, voice, and video.',
        'This expands the reachable workflow base from knowledge lookup into richer operational work.',
        'Multimodality matters most in service, field operations, claims, compliance, and training.'
    ], accent=BLUE, fill=LIGHT_BLUE)
    card(s, Inches(4.79), Inches(1.55), Inches(3.75), Inches(4.9), '5. From standalone chat to embedded workflow', [
        'The winning experiences disappear into CRM, ERP, ITSM, contact center, and internal operating tools.',
        'Users adopt agents faster when the agent works inside the system of record and respects process context.',
        'Integration depth is becoming more strategic than interface novelty.'
    ], accent=GREEN, fill=LIGHT_GREEN)
    card(s, Inches(8.78), Inches(1.55), Inches(3.75), Inches(4.9), '6. From pilot excitement to ROI and control', [
        'Management attention is shifting from proof-of-concept demos to economics, reliability, and risk management.',
        'That elevates evaluation, permissions, human escalation, observability, and audit trails.',
        'Governance is no longer a brake on adoption; it is the enabler of adoption at scale.'
    ], accent=RED, fill=LIGHT_RED)
    add_source(s, COMMON_SOURCE, num)


def slide_implications(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'What leading enterprises are building now | A digital employee operating model, not another AI point solution')
    add_rect(s, LM, Inches(1.42), CONTENT_W, Inches(0.8), NAVY)
    add_text(s, Inches(1.0), Inches(1.52), Inches(11.0), Inches(0.56),
             'Scaling digital employees requires simultaneous moves across workflow design, technical control, and organization design.',
             font_size=BODY_SIZE, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    pillars = [
        ('Workflow layer', ['Select high-volume work with clear triggers and measurable outputs', 'Define where agents act autonomously and where humans approve', 'Redesign exceptions, not just the happy path']),
        ('Platform layer', ['Connect agents to trusted data, tools, and permissions', 'Standardize evaluation, telemetry, prompt/version control, and routing', 'Treat agent identity and memory as enterprise architecture issues']),
        ('Organization layer', ['Create ownership across business, tech, risk, and operations', 'Train managers to supervise mixed human-digital teams', 'Move KPIs from activity counts toward throughput and quality outcomes']),
    ]
    pw = Inches(3.72)
    gap = Inches(0.29)
    fills = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_ORANGE]
    accents = [BLUE, GREEN, ORANGE]
    for i, ((title, items), fill, accent) in enumerate(zip(pillars, fills, accents)):
        x = LM + (pw + gap) * i
        add_rect(s, x, Inches(2.6), pw, Inches(3.55), fill)
        add_rect(s, x, Inches(2.6), pw, Inches(0.58), accent)
        add_text(s, x + Inches(0.12), Inches(2.6), pw - Inches(0.24), Inches(0.58), title,
                 font_size=BODY_SIZE, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        bullet_block(s, x + Inches(0.18), Inches(3.38), pw - Inches(0.36), items, Inches(2.45), size=Pt(12), spacing=Pt(9))
    add_rect(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.36), BG)
    add_text(s, Inches(1.0), Inches(6.4), Inches(11.2), Inches(0.2),
             'Bottom line: firms that deploy digital employees successfully behave like operating model designers, not just software buyers.',
             font_size=Pt(12), font_color=NAVY, bold=True)
    add_source(s, COMMON_SOURCE, num)


def slide_governance(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'Scaling safely | Capability without control will stall enterprise trust')
    left_x = Inches(0.8)
    right_x = Inches(6.75)
    col_w = Inches(5.55)
    add_text(s, left_x, Inches(1.55), col_w, Inches(0.35), 'Where digital employees fail',
             font_size=SUB_SIZE, font_color=NAVY, bold=True)
    add_hline(s, left_x, Inches(2.0), col_w, LINE)
    bullet_block(s, left_x, Inches(2.18), col_w, [
        'Hallucinated actions or incorrect reasoning inside sensitive workflows',
        'Tool misuse caused by excessive privileges or poor identity design',
        'Process drift when prompts, rules, or data sources change without control',
        'Low adoption when users cannot understand what the agent did or why'
    ], Inches(3.15), size=Pt(12))
    add_text(s, right_x, Inches(1.55), col_w, Inches(0.35), 'What enterprises need to manage',
             font_size=SUB_SIZE, font_color=NAVY, bold=True)
    add_hline(s, right_x, Inches(2.0), col_w, LINE)
    bullet_block(s, right_x, Inches(2.18), col_w, [
        'Role-based permissions, sandboxing, and clear human escalation points',
        'Evaluation loops covering quality, safety, latency, and business outcomes',
        'Traceability: logs, audit trails, version control, and incident review',
        'Change management so managers know how to supervise digital work'
    ], Inches(3.15), size=Pt(12))
    add_rect(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.75), NAVY)
    add_text(s, Inches(1.0), Inches(6.0), Inches(11.2), Inches(0.5),
             'Strategic implication: in the next phase of the market, trust architecture may matter more than model architecture.',
             font_size=BODY_SIZE, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_source(s, COMMON_SOURCE, num)


def slide_roadmap(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, 'A pragmatic 12–18 month roadmap | Start with constrained work, then industrialize')
    add_hline(s, Inches(1.3), Inches(3.25), Inches(10.4), LINE, Pt(2))
    phases = [
        ('1', 'Prioritize', 'Pick 2–3 roles with high volume, clear SOPs, and painful capacity gaps.'),
        ('2', 'Pilot', 'Launch narrow use cases with explicit success metrics, approvals, and fallback paths.'),
        ('3', 'Industrialize', 'Add connectors, evaluation, observability, identity controls, and support processes.'),
        ('4', 'Redesign', 'Rewrite team charters, KPIs, and manager routines around mixed human-digital delivery.'),
    ]
    xs = [Inches(1.1), Inches(4.0), Inches(6.9), Inches(9.8)]
    for (idx, title, desc), x in zip(phases, xs):
        add_oval(s, x, Inches(3.0), idx, size=Inches(0.5), bg=NAVY)
        add_text(s, x - Inches(0.2), Inches(2.25), Inches(0.9), Inches(0.4), title,
                 font_size=BODY_SIZE, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
        add_rect(s, x - Inches(1.0), Inches(3.75), Inches(2.1), Inches(1.6), BG)
        add_text(s, x - Inches(0.88), Inches(3.95), Inches(1.86), Inches(1.15),
                 desc, font_size=Pt(11.5), font_color=DARK, alignment=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.65), BG)
    add_text(s, Inches(1.0), Inches(6.03), Inches(11.2), Inches(0.42),
             'Do not begin with the most powerful agent. Begin with the most governable workflow that can produce visible business value.',
             font_size=BODY_SIZE, font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_source(s, COMMON_SOURCE, num)


def slide_closing(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
    add_text(s, Inches(1.2), Inches(2.05), Inches(10.9), Inches(1.0),
             'Digital employees will not simply automate work.',
             font_size=Pt(26), font_name='Georgia', font_color=NAVY,
             bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.2), Inches(3.0), Inches(10.9), Inches(0.9),
             'They will reshape how work is designed, governed, and measured.',
             font_size=SUB_SIZE, font_color=DARK, alignment=PP_ALIGN.CENTER)
    add_hline(s, Inches(5.45), Inches(4.05), Inches(2.45), NAVY, Pt(1.5))
    add_text(s, Inches(1.6), Inches(4.45), Inches(10.2), Inches(1.05), [
             'The immediate challenge for management is not whether to use them,',
             'but where to trust them first.'
             ], font_size=BODY_SIZE, font_color=MED, alignment=PP_ALIGN.CENTER, paragraph_space=Pt(8))
    add_hline(s, Inches(1.0), Inches(6.8), Inches(3.0), NAVY, Pt(2))
    add_source(s, COMMON_SOURCE, num)


def build(path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide_cover(prs)
    slide_exec_summary(prs, 2)
    slide_why_now(prs, 3)
    slide_trends_a(prs, 4)
    slide_trends_b(prs, 5)
    slide_implications(prs, 6)
    slide_governance(prs, 7)
    slide_roadmap(prs, 8)
    slide_closing(prs, 9)
    prs.save(path)
    full_cleanup(path)


if __name__ == '__main__':
    out_dir = '/Users/kaku/WorkBuddy/Claw/output'
    os.makedirs(out_dir, exist_ok=True)
    out_file = os.path.join(out_dir, 'digital_employee_trends_mckinsey_style.pptx')
    build(out_file)
    print(out_file)
