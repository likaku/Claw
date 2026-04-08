from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os
import zipfile

SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.5)
SOURCE_Y = Inches(7.1)

NAVY = RGBColor(0x05, 0x1C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(0x33, 0x33, 0x33)
MED = RGBColor(0x66, 0x66, 0x66)
LINE = RGBColor(0xCC, 0xCC, 0xCC)
BG = RGBColor(0xF2, 0xF2, 0xF2)
BLUE = RGBColor(0x00, 0x6B, 0xA6)
GREEN = RGBColor(0x00, 0x7A, 0x53)
ORANGE = RGBColor(0xD4, 0x6A, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
RED = RGBColor(0xC6, 0x28, 0x28)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)

TITLE_SIZE = Pt(18)
SUB_SIZE = Pt(13)
BODY_SIZE = Pt(10)
SMALL_SIZE = Pt(8)


def clean_shape(shape):
    style = shape._element.find(qn('p:style'))
    if style is not None:
        shape._element.remove(style)


def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def add_text(slide, left, top, width, height, text, font_size=BODY_SIZE,
             font_name='Arial', font_color=DARK, bold=False,
             alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             paragraph_space=Pt(4), line_multiple=1.25):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '36000')
    lines = text if isinstance(text, list) else [text]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = paragraph_space if idx > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * line_multiple)
        for run in p.runs:
            set_ea_font(run)
    return box


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    clean_shape(shape)
    return shape


def add_hline(slide, left, top, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, left, top, length, h, color)


def add_oval(slide, left, top, label, size=Inches(0.36), bg=NAVY, fg=WHITE, font_size=Pt(9)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    clean_shape(shape)
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = font_size
    p.font.name = 'Arial'
    p.font.color.rgb = fg
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    bodyPr = shape.text_frame._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    for run in p.runs:
        set_ea_font(run)
    return shape


def add_source(slide, text):
    add_text(slide, Inches(0.5), SOURCE_Y, Inches(11.4), Inches(0.18), text,
             font_size=SMALL_SIZE, font_color=MED)
    add_text(slide, Inches(12.3), SOURCE_Y, Inches(0.5), Inches(0.18), '1/1',
             font_size=SMALL_SIZE, font_color=MED, alignment=PP_ALIGN.RIGHT)


def full_cleanup(path):
    tmp = path + '.tmp'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith('.xml'):
                root = etree.fromstring(data)
                for style in root.findall(f'.//{{{ns_p}}}style'):
                    style.getparent().remove(style)
                if 'theme' in item.filename.lower():
                    for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                        for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                            el.getparent().remove(el)
                data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes')
            zout.writestr(item, data)
    os.replace(tmp, path)


def build(path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # === 标题 ===
    add_text(s, LM, Inches(0.12), Inches(12.2), Inches(0.55),
             'OpenClaw 创始过程｜Peter Steinberger 如何把个人 AI 助手做成现象级项目',
             font_size=TITLE_SIZE, font_name='Georgia', font_color=BLACK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_hline(s, LM, Inches(0.62), Inches(12.2), NAVY, Pt(1))

    # ==================== ROW 1: 创始人档案 + 核心洞察 ====================
    # 左侧：创始人档案卡片
    add_rect(s, Inches(0.5), Inches(0.78), Inches(4.0), Inches(2.85), BG)
    add_oval(s, Inches(0.68), Inches(1.0), 'PS', size=Inches(0.42), bg=NAVY)
    add_text(s, Inches(1.2), Inches(1.02), Inches(2.8), Inches(0.28),
             'Peter Steinberger', font_size=Pt(14), font_color=NAVY, bold=True)
    add_text(s, Inches(1.2), Inches(1.32), Inches(2.8), Inches(0.22),
             '奥地利开发者 / 开源作者 / OpenClaw 发起人', font_size=Pt(9), font_color=MED)
    
    # 分隔线
    add_hline(s, Inches(0.68), Inches(1.65), Inches(3.6), LINE)
    
    # 创始人关键标签
    labels = [
        ('背景', 'iOS/Swift 生态资深开发者', LIGHT_BLUE, BLUE),
        ('经历', 'PSPDFKit 创业 → AI 工具实验', LIGHT_GREEN, GREEN),
        ('转向', '2026.02 加入 OpenAI', LIGHT_ORANGE, ORANGE),
        ('理念', '改变世界 > 做大公司', LIGHT_PURPLE, PURPLE),
    ]
    y_start = Inches(1.78)
    for i, (icon, text, fill, accent) in enumerate(labels):
        y = y_start + i * Inches(0.42)
        add_rect(s, Inches(0.68), y, Inches(3.6), Inches(0.36), fill)
        add_rect(s, Inches(0.68), y, Inches(0.06), Inches(0.36), accent)
        add_text(s, Inches(0.84), y + Inches(0.06), Inches(0.6), Inches(0.24),
                 icon, font_size=Pt(9), font_color=NAVY, bold=True)
        add_text(s, Inches(1.5), y + Inches(0.06), Inches(2.6), Inches(0.24),
                 text, font_size=Pt(9), font_color=DARK)

    # 右侧：核心洞察（Big Number 风格）
    add_rect(s, Inches(4.65), Inches(0.78), Inches(3.95), Inches(2.85), LIGHT_BLUE)
    add_rect(s, Inches(4.65), Inches(0.78), Inches(0.1), Inches(2.85), BLUE)
    add_text(s, Inches(4.9), Inches(0.95), Inches(3.5), Inches(0.26),
             '创始洞察', font_size=Pt(11), font_color=BLUE, bold=True)
    add_text(s, Inches(4.9), Inches(1.35), Inches(3.5), Inches(1.1),
             '"不是再做会聊天的AI，\n而是能真正干活的个人助手"',
             font_size=Pt(13), font_color=NAVY, bold=True, line_multiple=1.2)
    
    # 3个关键特征
    features = [
        ('本地优先', '数据隐私 + 离线可用'),
        ('跨平台', '多端统一体验'),
        ('工具调用', '连接真实系统'),
    ]
    for i, (title, desc) in enumerate(features):
        x = Inches(4.9) + i * Inches(1.15)
        add_oval(s, x, Inches(2.55), str(i+1), size=Inches(0.28), bg=BLUE, font_size=Pt(8))
        add_text(s, x + Inches(0.35), Inches(2.55), Inches(0.75), Inches(0.2),
                 title, font_size=Pt(9), font_color=NAVY, bold=True)
        add_text(s, x, Inches(2.85), Inches(1.0), Inches(0.4),
                 desc, font_size=Pt(8), font_color=MED)

    # ==================== ROW 2: 项目演进 + 技术架构 ====================
    # 项目演进时间线
    add_text(s, Inches(0.5), Inches(3.72), Inches(4.0), Inches(0.26),
             '项目演进路径', font_size=Pt(11), font_color=NAVY, bold=True)
    
    timeline = [
        ('1', '前史', '移动开发→AI工具', LIGHT_ORANGE, ORANGE),
        ('2', 'Clawdbot', '"能干活的助手"原型', LIGHT_BLUE, BLUE),
        ('3', 'Moltbot', '品牌探索期', LIGHT_GREEN, GREEN),
        ('4', 'OpenClaw', '定位稳定+开源独立', LIGHT_PURPLE, PURPLE),
    ]
    for i, (num, title, desc, fill, accent) in enumerate(timeline):
        x = Inches(0.5) + i * Inches(0.98)
        add_oval(s, x + Inches(0.32), Inches(4.05), num, size=Inches(0.28), bg=accent, font_size=Pt(8))
        if i < len(timeline) - 1:
            add_hline(s, x + Inches(0.65), Inches(4.17), Inches(0.65), LINE, Pt(2))
        add_text(s, x, Inches(4.42), Inches(0.95), Inches(0.22),
                 title, font_size=Pt(9), font_color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s, x, Inches(4.68), Inches(0.95), Inches(0.45),
                 desc, font_size=Pt(8), font_color=MED, alignment=PP_ALIGN.CENTER, line_multiple=1.15)

    # 技术架构卡片
    add_rect(s, Inches(4.65), Inches(3.72), Inches(3.95), Inches(1.55), BG)
    add_text(s, Inches(4.85), Inches(3.88), Inches(3.5), Inches(0.22),
             '技术架构特征', font_size=Pt(11), font_color=NAVY, bold=True)
    
    arch_items = [
        ('本地优先', '数据不出设备，隐私可控'),
        ('消息渠道', '接入 Telegram/微信等主流平台'),
        ('工具生态', '调用 API、文件系统、自动化脚本'),
        ('开源独立', '基金会治理，社区驱动'),
    ]
    for i, (title, desc) in enumerate(arch_items):
        y = Inches(4.18) + i * Inches(0.26)
        add_text(s, Inches(4.85), y, Inches(1.0), Inches(0.22),
                 f'• {title}', font_size=Pt(8.5), font_color=NAVY, bold=True)
        add_text(s, Inches(5.9), y, Inches(2.5), Inches(0.22),
                 desc, font_size=Pt(8.5), font_color=DARK)

    # ==================== ROW 3: 战略转向 + 市场定位 ====================
    # 战略转向卡片
    add_rect(s, Inches(0.5), Inches(5.35), Inches(4.0), Inches(1.6), LIGHT_ORANGE)
    add_rect(s, Inches(0.5), Inches(5.35), Inches(0.1), Inches(1.6), ORANGE)
    add_text(s, Inches(0.75), Inches(5.52), Inches(3.6), Inches(0.24),
             '2026年战略转向', font_size=Pt(11), font_color=ORANGE, bold=True)
    
    # 两个并列模块
    add_rect(s, Inches(0.75), Inches(5.85), Inches(1.75), Inches(0.95), WHITE)
    add_text(s, Inches(0.88), Inches(5.95), Inches(1.5), Inches(0.2),
             '创始人去向', font_size=Pt(9), font_color=NAVY, bold=True)
    add_text(s, Inches(0.88), Inches(6.18), Inches(1.5), Inches(0.55),
             '加入 OpenAI\n投入"下一代个人agent"', font_size=Pt(8.5), font_color=DARK, line_multiple=1.15)
    
    add_rect(s, Inches(2.6), Inches(5.85), Inches(1.75), Inches(0.95), WHITE)
    add_text(s, Inches(2.73), Inches(5.95), Inches(1.5), Inches(0.2),
             'OpenClaw去向', font_size=Pt(9), font_color=NAVY, bold=True)
    add_text(s, Inches(2.73), Inches(6.18), Inches(1.5), Inches(0.55),
             '进入基金会\n保持开源独立运营', font_size=Pt(8.5), font_color=DARK, line_multiple=1.15)

    # 市场定位对比表
    add_text(s, Inches(4.65), Inches(5.35), Inches(3.95), Inches(0.24),
             '市场定位对比', font_size=Pt(11), font_color=NAVY, bold=True)
    
    # 表头
    add_rect(s, Inches(4.65), Inches(5.65), Inches(3.95), Inches(0.28), NAVY)
    add_text(s, Inches(4.75), Inches(5.68), Inches(1.2), Inches(0.2),
             '维度', font_size=Pt(8), font_color=WHITE, bold=True)
    add_text(s, Inches(6.0), Inches(5.68), Inches(1.25), Inches(0.2),
             'OpenClaw', font_size=Pt(8), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(7.3), Inches(5.68), Inches(1.2), Inches(0.2),
             '传统AI助手', font_size=Pt(8), font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 表格内容
    table_data = [
        ('数据隐私', '本地优先', '云端处理'),
        ('执行能力', '真实工具调用', '对话为主'),
        ('商业模式', '开源+基金会', 'SaaS订阅'),
        ('生态开放', '社区驱动', '平台封闭'),
    ]
    for i, (dim, v1, v2) in enumerate(table_data):
        y = Inches(5.98) + i * Inches(0.24)
        fill = BG if i % 2 == 0 else WHITE
        add_rect(s, Inches(4.65), y, Inches(3.95), Inches(0.24), fill)
        add_text(s, Inches(4.75), y + Inches(0.02), Inches(1.2), Inches(0.2),
                 dim, font_size=Pt(8), font_color=NAVY, bold=True)
        add_text(s, Inches(6.0), y + Inches(0.02), Inches(1.25), Inches(0.2),
                 v1, font_size=Pt(8), font_color=GREEN, alignment=PP_ALIGN.CENTER)
        add_text(s, Inches(7.3), y + Inches(0.02), Inches(1.2), Inches(0.2),
                 v2, font_size=Pt(8), font_color=MED, alignment=PP_ALIGN.CENTER)

    # ==================== 右侧：关键里程碑数字 ====================
    add_rect(s, Inches(8.7), Inches(0.78), Inches(4.05), Inches(2.85), NAVY)
    add_text(s, Inches(8.9), Inches(0.95), Inches(3.6), Inches(0.24),
             '关键里程碑', font_size=Pt(11), font_color=WHITE, bold=True)
    
    milestones = [
        ('GitHub Stars', '50K+', '开源社区关注度'),
        ('用户覆盖', '100+', '国家和地区'),
        ('集成渠道', '15+', '消息平台/工具'),
        ('贡献者', '200+', '全球开发者'),
    ]
    for i, (label, num, desc) in enumerate(milestones):
        y = Inches(1.35) + i * Inches(0.55)
        add_text(s, Inches(8.9), y, Inches(1.4), Inches(0.2),
                 label, font_size=Pt(9), font_color=MED)
        add_text(s, Inches(10.4), y, Inches(1.0), Inches(0.26),
                 num, font_size=Pt(14), font_color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
        add_text(s, Inches(8.9), y + Inches(0.22), Inches(2.5), Inches(0.18),
                 desc, font_size=Pt(8), font_color=RGBColor(0x88, 0xAA, 0xCC))

    # 行业影响力
    add_rect(s, Inches(8.7), Inches(3.72), Inches(4.05), Inches(1.55), BG)
    add_text(s, Inches(8.9), Inches(3.88), Inches(3.6), Inches(0.22),
             '行业影响力', font_size=Pt(11), font_color=NAVY, bold=True)
    
    impacts = [
        '推动"本地优先AI"成为行业共识',
        '验证个人AI助手的产品形态',
        '为开源AI项目提供治理范式',
        '引领"agent工具化"技术路线',
    ]
    for i, text in enumerate(impacts):
        y = Inches(4.18) + i * Inches(0.26)
        add_oval(s, Inches(8.9), y, str(i+1), size=Inches(0.22), bg=BLUE, font_size=Pt(7))
        add_text(s, Inches(9.2), y, Inches(3.3), Inches(0.22),
                 text, font_size=Pt(8.5), font_color=DARK)

    # 一句话总结
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.25), Inches(0.35), BG)
    add_text(s, Inches(0.7), Inches(6.58), Inches(11.8), Inches(0.28),
             '核心结论：OpenClaw 不是从"大公司产品路线"长出来的，而是从创始人对"真正能做事的个人 AI"这一强烈产品判断中长出来的。',
             font_size=Pt(10), font_color=NAVY, bold=True)

    add_source(s, 'Source: OpenClaw GitHub / Peter Steinberger personal site / TechCrunch (2026-02-15) / Lex Fridman episode page')

    prs.save(path)
    full_cleanup(path)


if __name__ == '__main__':
    out_dir = '/Users/kaku/WorkBuddy/Claw/output'
    os.makedirs(out_dir, exist_ok=True)
    out_file = os.path.join(out_dir, 'openclaw_founder_onepager.pptx')
    build(out_file)
    print(out_file)
