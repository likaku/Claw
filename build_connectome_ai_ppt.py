from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os
import zipfile


SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.8)
RM = Inches(0.8)
CONTENT_W = SW - LM - RM
TITLE_Y = Inches(0.15)
TITLE_H = Inches(0.9)
CONTENT_TOP = Inches(1.35)
SOURCE_Y = Inches(7.03)

NAVY = RGBColor(0x05, 0x1C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
MED = RGBColor(0x66, 0x66, 0x66)
LINE = RGBColor(0xCC, 0xCC, 0xCC)
BG = RGBColor(0xF2, 0xF2, 0xF2)
BLUE = RGBColor(0x00, 0xA9, 0xF4)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BLUE = RGBColor(0xE9, 0xF6, 0xFD)
LIGHT_GREEN = RGBColor(0xE9, 0xF5, 0xEE)
LIGHT_ORANGE = RGBColor(0xFD, 0xF2, 0xE9)
LIGHT_RED = RGBColor(0xFB, 0xEB, 0xEB)

TITLE_SIZE = Pt(22)
SUB_SIZE = Pt(18)
BODY_SIZE = Pt(14)
SMALL_SIZE = Pt(9)
BIG_SIZE = Pt(32)


def clean_shape(shape):
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def add_text(slide, left, top, width, height, text, font_size=BODY_SIZE,
             font_name='Arial', font_color=DARK, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP,
             paragraph_space=Pt(6), line_multiple=1.35):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = paragraph_space if idx > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * line_multiple)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return box


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    clean_shape(shape)
    return shape


def add_hline(slide, left, top, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, left, top, length, h, color)


def add_oval(slide, left, top, label, size=Inches(0.46), bg=NAVY, fg=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    clean_shape(shape)
    tf = shape.text_frame
    tf.paragraphs[0].text = label
    p = tf.paragraphs[0]
    p.font.size = BODY_SIZE
    p.font.name = 'Arial'
    p.font.color.rgb = fg
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    for run in p.runs:
        set_ea_font(run, 'KaiTi')
    return shape


def add_action_title(slide, text):
    add_text(slide, LM, TITLE_Y, CONTENT_W, TITLE_H, text,
             font_size=TITLE_SIZE, font_name='Georgia', font_color=BLACK,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, LM, Inches(1.03), CONTENT_W, BLACK, Pt(0.5))


def add_source(slide, text):
    add_text(slide, LM, SOURCE_Y, Inches(11.2), Inches(0.22), text,
             font_size=SMALL_SIZE, font_name='Arial', font_color=MED)


def add_page_number(slide, num):
    add_text(slide, Inches(12.3), SOURCE_Y, Inches(0.35), Inches(0.22), str(num),
             font_size=SMALL_SIZE, font_name='Arial', font_color=MED,
             alignment=PP_ALIGN.RIGHT)


def full_cleanup(path):
    tmp = path + '.tmp'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    with zipfile.ZipFile(path, 'r') as zin, zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith('.xml'):
                root = etree.fromstring(data)
                for style in root.findall(f'.//{{{ns_p}}}style'):
                    style.getparent().remove(style)
                if 'theme' in item.filename.lower():
                    for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                        for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                            el.getparent().remove(el)
                data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes')
            zout.writestr(item, data)
    os.replace(tmp, path)


def bullet_block(slide, x, y, w, items, color=DARK, size=BODY_SIZE, spacing=Pt(10), height=Inches(3.8)):
    text = [f'• {item}' for item in items]
    add_text(slide, x, y, w, height, text, font_size=size, font_color=color, paragraph_space=spacing)


def card(slide, x, y, w, h, title, bullets, accent=NAVY, fill=BG):
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, w, Inches(0.55), accent)
    add_text(slide, x + Inches(0.12), y, w - Inches(0.24), Inches(0.55), title,
             font_size=BODY_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    bullet_block(slide, x + Inches(0.18), y + Inches(0.75), w - Inches(0.36), bullets,
                 color=DARK, size=BODY_SIZE, spacing=Pt(8), height=h - Inches(0.95))


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
    add_text(s, Inches(1.0), Inches(1.7), Inches(11.2), Inches(1.4),
             'Mck-ppt-design-skill 示例演示：\n连接组学与 AI 驱动的大脑动态建模',
             font_size=Pt(28), font_name='Georgia', font_color=NAVY,
             bold=True, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_multiple=1.2)
    add_text(s, Inches(1.4), Inches(3.55), Inches(10.5), Inches(0.8),
             '为什么“只有接线图还不够”，以及为什么 AI 正成为预测全脑活动的关键工具',
             font_size=Pt(18), font_color=DARK, alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(4.55), Inches(11.2), Inches(0.6),
             '基于用户提供文献摘要整理｜版式按 Mck-ppt-design-skill 规范生成｜2026 年 3 月',
             font_size=BODY_SIZE, font_color=MED, alignment=PP_ALIGN.CENTER)
    add_hline(s, Inches(1.0), Inches(6.75), Inches(3.2), NAVY, Pt(2))
    add_text(s, Inches(1.0), Inches(6.9), Inches(11.0), Inches(0.25),
             '作者/生成方式：GPT + python-pptx + Mck-ppt-design-skill',
             font_size=SMALL_SIZE, font_color=MED)
    return s


def slide_skill_intro(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '为何用 Mck-ppt-design-skill 来表达这类跨学科综述')
    items = [
        ('结构化叙事', '把“背景-论点-证据-结论”组织成咨询式页面，便于从复杂文献中提炼主张。'),
        ('统一视觉层级', '固定标题、正文、脚注与强调色体系，让跨学科信息在一页内也保持清晰可读。'),
        ('稳定工程化交付', '通过中文字体控制、行距控制与 XML 清理，降低 PPT 打不开、错位、重叠等问题。'),
    ]
    ys = [Inches(1.55), Inches(3.0), Inches(4.45)]
    fills = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_ORANGE]
    accents = [BLUE, GREEN, ORANGE]
    for i, ((title, desc), y, fill, accent) in enumerate(zip(items, ys, fills, accents), 1):
        add_oval(s, LM, y, str(i))
        add_text(s, LM + Inches(0.65), y - Inches(0.02), Inches(3.0), Inches(0.35),
                 title, font_size=SUB_SIZE, font_color=NAVY, bold=True)
        add_rect(s, Inches(4.0), y - Inches(0.05), Inches(8.55), Inches(0.95), fill)
        add_text(s, Inches(4.22), y + Inches(0.08), Inches(8.1), Inches(0.65),
                 desc, font_size=BODY_SIZE, font_color=DARK)
        add_rect(s, Inches(4.0), y - Inches(0.05), Inches(0.12), Inches(0.95), accent)
    add_rect(s, Inches(0.8), Inches(5.95), CONTENT_W, Inches(0.65), BG)
    add_text(s, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.45),
             '本 deck 因此采用“先给结论，再给案例”的表达结构：更像研究负责人汇报，而不是线性读论文。',
             font_size=BODY_SIZE, font_color=NAVY, bold=True)
    add_source(s, 'Source: Mck-ppt-design-skill README / design specification, GitHub repository by likaku')
    add_page_number(s, num)
    return s


def slide_exec_summary(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '执行摘要｜该领域的四个关键判断')
    add_rect(s, LM, Inches(1.42), CONTENT_W, Inches(0.9), NAVY)
    add_text(s, Inches(1.05), Inches(1.52), Inches(10.9), Inches(0.7),
             '核心结论：连接组学提供结构边界，AI 与动态模型提供功能解释；两者结合，才可能走向可信的全脑预测。',
             font_size=SUB_SIZE, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    cols = [
        ('1', '静态连接图不够', '没有动力学、电突触、神经调节与状态依赖信息，无法单独解释脑功能。'),
        ('2', '连接组约束 AI 更可解释', '真实稀疏连接减少自由度，更容易逼近生物机制，并做出单细胞级预测。'),
        ('3', 'AI 已能预测全脑活动', 'RNN、Transformer 与基础模型正从拟合历史走向预测未来脑状态。'),
        ('4', '统一基准决定领域增速', '像 WeatherBench 一样，神经科学需要公开基准来比较模型、数据和训练策略。'),
    ]
    cw = Inches(2.72)
    gap = Inches(0.2)
    for i, (idx, title, desc) in enumerate(cols):
        x = LM + (cw + gap) * i
        add_rect(s, x, Inches(2.7), cw, Inches(3.15), BG)
        add_oval(s, x + Inches(0.12), Inches(2.82), idx)
        add_text(s, x + Inches(0.65), Inches(2.82), cw - Inches(0.8), Inches(0.4),
                 title, font_size=BODY_SIZE, font_color=NAVY, bold=True)
        add_hline(s, x + Inches(0.15), Inches(3.33), cw - Inches(0.3), LINE)
        add_text(s, x + Inches(0.15), Inches(3.5), cw - Inches(0.3), Inches(2.1),
                 desc, font_size=BODY_SIZE, font_color=DARK)
    add_source(s, 'Source: Based on user-provided literature synthesis across worm, fly, zebrafish, and mouse studies')
    add_page_number(s, num)
    return s


def slide_connectome_limit(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '观点一｜仅有连接组学是巨大的解剖成就，但仍不足以单独解释脑功能')
    add_text(s, LM, Inches(1.48), Inches(7.6), Inches(0.38),
             '“接线图”告诉我们谁连到谁；但真正决定输出的，是连接如何在时间中被激活、调制与重写。',
             font_size=SUB_SIZE, font_color=NAVY, bold=True)
    add_hline(s, LM, Inches(1.95), Inches(7.6), LINE)
    bullet_block(s, LM, Inches(2.15), Inches(7.55), [
        '大规模连接组绘制（如果蝇全脑、人脑皮层片段）极大提升了我们对结构约束的认识。',
        '但静态图谱缺少神经元内在动力学、短时程/长时程可塑性、电突触与神经调质状态。',
        '同一结构在不同任务、觉醒水平或感觉背景下，可表现出完全不同的有效连接与动力学。',
        '因此真正需要回答的问题，不是“线路图是什么”，而是“这张图在什么规则下运行”。',
    ], height=Inches(3.9))
    add_rect(s, Inches(9.0), Inches(1.5), Inches(3.5), Inches(5.0), BG)
    add_text(s, Inches(9.22), Inches(1.72), Inches(3.05), Inches(0.32),
             'Key Takeaways', font_size=BODY_SIZE, font_color=NAVY, bold=True)
    add_hline(s, Inches(9.22), Inches(2.18), Inches(3.0), LINE)
    add_text(s, Inches(9.22), Inches(2.35), Inches(3.05), Inches(3.9), [
        '1. 结构是必要约束，不是充分解释。',
        '2. 功能需要时间维与状态维。',
        '3. 计算模型负责填补“结构→功能”的空白。',
    ], font_size=BODY_SIZE, font_color=DARK, paragraph_space=Pt(10))
    add_source(s, 'Source: Connectomics + functional modeling argument summarized from the provided research narrative')
    add_page_number(s, num)
    return s


def slide_connectome_ai(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '观点二｜连接组约束的 AI 模型更容易逼近真实生物计算机制')
    add_rect(s, LM, Inches(1.46), CONTENT_W, Inches(0.78), NAVY)
    add_text(s, Inches(1.0), Inches(1.56), Inches(11.0), Inches(0.56),
             '真实稀疏连接将搜索空间从“任何网络都可以”收缩为“只有生物上可能的网络可以”——可解释性因此显著提高。',
             font_size=BODY_SIZE, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    card(s, Inches(0.8), Inches(2.55), Inches(3.8), Inches(3.35), '机制层收益', [
        '显著减少自由参数与等价解。',
        '让训练更聚焦于单神经元与突触参数。',
        '更容易把“任务成功”映射回“生物机制成立”。',
    ], accent=BLUE, fill=LIGHT_BLUE)
    card(s, Inches(4.77), Inches(2.55), Inches(3.8), Inches(3.35), '预测层收益', [
        '可做到单细胞活动预测，而不只是群体平均。',
        '更适合与独立实验数据逐项比对。',
        '更利于发现哪些参数是关键瓶颈。',
    ], accent=GREEN, fill=LIGHT_GREEN)
    card(s, Inches(8.74), Inches(2.55), Inches(3.8), Inches(3.35), '果蝇案例', [
        '视觉运动通路模型覆盖 64 类细胞。',
        '连接方式来自实验测量的真实连接组。',
        '训练后既完成任务，也匹配 26 项独立细胞响应测量。',
    ], accent=ORANGE, fill=LIGHT_ORANGE)
    add_source(s, 'Source: Drosophila connectome-constrained modeling examples; see also Nature/Nature Neuroscience search results on connectome-constrained recurrent networks')
    add_page_number(s, num)
    return s


def slide_prediction(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '观点三｜人工智能已成为预测与模拟全脑神经活动的强大工具')
    add_text(s, LM, Inches(1.48), Inches(11.6), Inches(0.42),
             '从线虫到小鼠，模型能力正在从“便宜地复现已有动力学”扩展到“跨个体、跨刺激、跨时间地预测未来活动”。',
             font_size=SUB_SIZE, font_color=NAVY, bold=True)
    cw = Inches(5.55)
    gap = Inches(0.75)
    left = LM
    right = LM + cw + gap
    add_text(s, left, Inches(2.0), cw, Inches(0.38), 'A. 斑马鱼：直接吃 3D 体视频，保留空间关系',
             font_size=BODY_SIZE, font_color=NAVY, bold=True)
    add_hline(s, left, Inches(2.46), cw, LINE)
    bullet_block(s, left, Inches(2.65), cw, [
        'ZAPBench 将全脑光片记录组织成统一预测基准。',
        '把 volumetric video 直接输入 3D U-Net，可隐式学习神经元空间邻近关系。',
        '其表现优于仅基于 1D trace 的时间序列模型。',
        'POCO 的“群体条件化”说明一个模型可在不同个体间快速适配。',
    ], height=Inches(3.6))
    add_text(s, right, Inches(2.0), cw, Inches(0.38), 'B. 小鼠：基础模型开始具备迁移与泛化能力',
             font_size=BODY_SIZE, font_color=NAVY, bold=True)
    add_hline(s, right, Inches(2.46), cw, LINE)
    bullet_block(s, right, Inches(2.65), cw, [
        'QuantFormer、BrainLM 利用海量双光子钙成像训练统一表示。',
        '掩码自编码与向量量化帮助模型在未见刺激上仍保持有效编码。',
        '只需少量新数据即可迁移到新个体，说明模型已具“脑活动基础模型”雏形。',
    ], height=Inches(3.6))
    add_source(s, 'Source: ZAPBench / POCO / QuantFormer / BrainLM examples summarized from the provided literature outline')
    add_page_number(s, num)
    return s


def slide_benchmark(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '观点四｜统一基准是推动领域从“好故事”走向“可累积进步”的关键')
    add_rect(s, Inches(0.8), Inches(1.55), Inches(3.4), Inches(1.85), NAVY)
    add_text(s, Inches(1.05), Inches(1.78), Inches(2.9), Inches(0.7),
             'WeatherBench\n给天气预测带来统一标尺',
             font_size=Pt(20), font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, Inches(4.7), Inches(1.8), Inches(7.9), Inches(1.2),
             '神经科学正在进入类似阶段：如果没有统一任务、统一数据拆分、统一评价指标，任何模型“有效”都难以横向比较。',
             font_size=BODY_SIZE, font_color=DARK)
    add_rect(s, LM, Inches(4.2), CONTENT_W, Inches(2.0), BG)
    labels = [
        ('统一任务', '如未来 1 步 / 多步预测、跨个体泛化、缺失神经元补全'),
        ('统一数据', '公开预处理、时间切分、训练/验证/测试划分'),
        ('统一指标', '相关系数、解释方差、多步滚动误差、跨个体迁移表现'),
        ('统一基线', 'trace-based / volumetric / RNN / Transformer / foundation model 并列对照'),
    ]
    xs = [Inches(1.0), Inches(3.95), Inches(6.9), Inches(9.85)]
    for i, ((title, desc), x) in enumerate(zip(labels, xs), 1):
        add_oval(s, x, Inches(4.45), str(i), size=Inches(0.42), bg=NAVY)
        add_text(s, x + Inches(0.5), Inches(4.42), Inches(2.1), Inches(0.32),
                 title, font_size=BODY_SIZE, font_color=NAVY, bold=True)
        add_text(s, x + Inches(0.5), Inches(4.84), Inches(2.25), Inches(1.0),
                 desc, font_size=Pt(12), font_color=DARK)
    add_source(s, 'Source: WeatherBench analogy and ZAPBench rationale from the user-provided summary; ZAPBench arXiv listing found via web search')
    add_page_number(s, num)
    return s


def slide_case_table(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '研究案例总览｜四类系统共同支撑“结构 + 动态 + AI”范式')
    headers = ['系统', '模型/方法', '证明了什么', '研究价值']
    widths = [Inches(1.6), Inches(3.1), Inches(4.25), Inches(2.75)]
    x = LM
    for h, w in zip(headers, widths):
        add_text(s, x, Inches(1.52), w, Inches(0.35), h,
                 font_size=BODY_SIZE, font_color=MED, bold=True)
        x += w
    add_hline(s, LM, Inches(1.98), CONTENT_W, BLACK, Pt(1.0))
    rows = [
        ('线虫', 'RNN / LSTM / GRU', '极小 GRU 也能高保真替代复杂生物物理模型', '把昂贵仿真压缩成可训练、可迭代代理模型'),
        ('果蝇', '连接组约束深度网络', '真实连接结构可支持运动检测并做出单细胞级预测', '把“任务表现”与“生物机制”绑定起来'),
        ('斑马鱼', 'ZAPBench + 3D U-Net + POCO', '直接利用 3D 时空视频与群体条件化可提升全脑预测', '推动跨个体、跨模型统一评测'),
        ('小鼠', 'QuantFormer / BrainLM', '基础模型能学到可迁移表示并泛化到新刺激/新个体', '提示神经活动领域也可能出现 foundation model'),
    ]
    row_h = Inches(1.08)
    y = Inches(2.15)
    for row in rows:
        x = LM
        for idx, (val, w) in enumerate(zip(row, widths)):
            add_text(s, x, y, w, row_h - Inches(0.12), val,
                     font_size=Pt(12) if idx > 0 else BODY_SIZE,
                     font_color=NAVY if idx == 0 else DARK,
                     bold=(idx == 0))
            x += w
        add_hline(s, LM, y + row_h, CONTENT_W, LINE, Pt(0.5))
        y += row_h
    add_source(s, 'Source: Worm / fly / zebrafish / mouse examples from the supplied summary and supplementary public paper searches')
    add_page_number(s, num)
    return s


def slide_implications(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(s, '结论与研究启示｜下一阶段的竞争力将来自“可解释预测”而非单点数据规模')
    add_rect(s, LM, Inches(1.45), CONTENT_W, Inches(0.9), NAVY)
    add_text(s, Inches(1.0), Inches(1.55), Inches(11.0), Inches(0.62),
             '下一代脑模型需要同时满足三件事：尊重真实结构、学习动态规律、并在统一基准上被严格比较。',
             font_size=SUB_SIZE, font_color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    card(s, Inches(0.8), Inches(2.75), Inches(3.72), Inches(3.0), '对实验神经科学', [
        '连接组与功能记录应同步设计。',
        '数据采集要为基准化与模型训练服务。',
        '独立实验验证仍是关键。',
    ], accent=BLUE, fill=LIGHT_BLUE)
    card(s, Inches(4.81), Inches(2.75), Inches(3.72), Inches(3.0), '对 AI 建模', [
        '少即是多：更强结构先验可换来更强解释力。',
        '多步预测与跨个体泛化将成为核心指标。',
        'foundation model 值得投入，但不能脱离生物约束。',
    ], accent=GREEN, fill=LIGHT_GREEN)
    card(s, Inches(8.82), Inches(2.75), Inches(3.72), Inches(3.0), '对领域组织方式', [
        '需要公开 benchmark、共享数据与标准报告。',
        '需要“模型—实验—基准”三方闭环。',
        '这会把神经科学推向更快的工程化迭代。',
    ], accent=RED, fill=LIGHT_RED)
    add_source(s, 'Source: Synthesized conclusion from the four arguments and the case studies in this deck')
    add_page_number(s, num)
    return s


def slide_closing(prs, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
    add_text(s, Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.1),
             '结构给出边界，动态给出规则，AI 给出预测能力。',
             font_size=Pt(26), font_name='Georgia', font_color=NAVY,
             bold=True, alignment=PP_ALIGN.CENTER)
    add_hline(s, Inches(5.45), Inches(3.25), Inches(2.45), NAVY, Pt(1.5))
    add_text(s, Inches(1.5), Inches(3.7), Inches(10.3), Inches(1.3),
             '谢谢。\n本文件由 Mck-ppt-design-skill 设计规范生成，可继续扩展为正式汇报版本。',
             font_size=SUB_SIZE, font_color=DARK, alignment=PP_ALIGN.CENTER,
             line_multiple=1.25)
    add_hline(s, Inches(1.0), Inches(6.8), Inches(3.0), NAVY, Pt(2))
    add_source(s, 'Source: Presentation generated in McKinsey-style design language using python-pptx helpers')
    add_page_number(s, num)
    return s


def build(path):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide_cover(prs)
    slide_skill_intro(prs, 2)
    slide_exec_summary(prs, 3)
    slide_connectome_limit(prs, 4)
    slide_connectome_ai(prs, 5)
    slide_prediction(prs, 6)
    slide_benchmark(prs, 7)
    slide_case_table(prs, 8)
    slide_implications(prs, 9)
    slide_closing(prs, 10)
    prs.save(path)
    full_cleanup(path)


if __name__ == '__main__':
    out_dir = '/Users/kaku/WorkBuddy/Claw/output'
    os.makedirs(out_dir, exist_ok=True)
    out_file = os.path.join(out_dir, 'connectome_ai_mck_skill_demo.pptx')
    build(out_file)
    print(out_file)
